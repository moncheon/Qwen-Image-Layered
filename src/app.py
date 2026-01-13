import gc
import torch

# ============================================================
# 설정 플래그
# ============================================================
# SageAttention 비활성화 (PyTorch 2.4에서 호환성 문제)
ENABLE_SAGE_ATTENTION = False

# 양자화 설정 플래그
# "4bit" - transformer+text_encoder 4bit (메모리 75% 감소, 품질 저하). 잡음만 나옴
# "4bit_quality" - transformer만 4bit (중간 품질, 권장). 잡음만 나옴
# "8bit" - 8bit 양자화 (VRAM 24GB 초과 가능)
# None - 양자화 없음 (bfloat16)
QUANTIZATION_MODE = None

# FirstBlockCache 설정
ENABLE_FIRST_BLOCK_CACHE = True
FIRST_BLOCK_CACHE_THRESHOLD = 0.1

# CPU Offload 설정 (THE GOD)
# "off" - GPU만 사용 (VRAM 많이 사용. VRAM 넘치면 임의로 RAM으로 넘김)
# "model" - 모델 단위 CPU offload (VRAM 넘치면 모델 단위로 고려하여 RAM으로 넘김)
# "sequential" - 레이어 단위 CPU offload (VRAM 넘치면 레이어 재료 단위로 고려하여 RAM으로 넘김)
CPU_OFFLOAD_MODE = "sequential"

# ============================================================
# 환경 검증 및 출력
# ============================================================
def print_environment_info():
    """검증된 환경 정보 출력"""
    print("=" * 60)
    print("환경 정보")
    print("=" * 60)

    # Python & PyTorch
    import sys
    print(f"Python: {sys.version.split()[0]}")
    print(f"PyTorch: {torch.__version__}")

    # CUDA
    if torch.cuda.is_available():
        print(f"CUDA: {torch.version.cuda}")
        print(f"GPU: {torch.cuda.get_device_name(0)}")
        vram_total = torch.cuda.get_device_properties(0).total_memory / (1024**3)
        print(f"VRAM: {vram_total:.1f} GB")
    else:
        print("CUDA: 사용 불가")

    # Diffusers
    try:
        import diffusers
        print(f"Diffusers: {diffusers.__version__}")
    except:
        print("Diffusers: 설치 안됨")

    # Transformers
    try:
        import transformers
        print(f"Transformers: {transformers.__version__}")
    except:
        print("Transformers: 설치 안됨")

    # BitsAndBytes
    try:
        import bitsandbytes
        print(f"BitsAndBytes: {bitsandbytes.__version__}")
        bnb_available = True
    except:
        print("BitsAndBytes: 설치 안됨")
        bnb_available = False

    # xformers
    try:
        import xformers
        print(f"xformers: {xformers.__version__}")
        xformers_available = True
    except:
        print("xformers: 설치 안됨")
        xformers_available = False

    # Triton (SageAttention 의존성)
    try:
        import triton
        print(f"Triton: {triton.__version__}")
        triton_available = True
    except:
        print("Triton: 설치 안됨")
        triton_available = False

    # SageAttention
    sageattention_available = False
    if triton_available:
        try:
            from sageattention import sageattn
            import sageattention
            ver = getattr(sageattention, '__version__', 'unknown')
            print(f"SageAttention: {ver}")
            sageattention_available = True
        except:
            print("SageAttention: 사용 불가")
    else:
        print("SageAttention: Triton 필요")

    print("-" * 60)
    print("SDPA 백엔드 상태")
    print("-" * 60)

    # Flash Attention (SDPA)
    flash_available = False
    mem_efficient_available = False
    try:
        flash_available = torch.backends.cuda.flash_sdp_enabled()
        print(f"Flash SDP: {'활성화' if flash_available else '비활성화'}")
    except:
        print("Flash SDP: 지원 안됨")

    try:
        mem_efficient_available = torch.backends.cuda.mem_efficient_sdp_enabled()
        print(f"Memory Efficient SDP: {'활성화' if mem_efficient_available else '비활성화'}")
    except:
        print("Memory Efficient SDP: 지원 안됨")

    # Math SDP (fallback)
    try:
        math_sdp = torch.backends.cuda.math_sdp_enabled()
        print(f"Math SDP (fallback): {'활성화' if math_sdp else '비활성화'}")
    except:
        pass

    print("-" * 60)
    print("현재 설정")
    print("-" * 60)
    print(f"QUANTIZATION_MODE: {QUANTIZATION_MODE}")
    print(f"ENABLE_SAGE_ATTENTION: {ENABLE_SAGE_ATTENTION}")
    print(f"ENABLE_FIRST_BLOCK_CACHE: {ENABLE_FIRST_BLOCK_CACHE}")
    if ENABLE_FIRST_BLOCK_CACHE:
        print(f"  - threshold: {FIRST_BLOCK_CACHE_THRESHOLD}")
    print(f"CPU_OFFLOAD_MODE: {CPU_OFFLOAD_MODE}")

    # 실제 사용될 Attention 백엔드 결정
    print("-" * 60)
    print("Attention 백엔드 (실제 사용)")
    print("-" * 60)
    if ENABLE_SAGE_ATTENTION and sageattention_available:
        print("→ SageAttention")
    elif xformers_available:
        print("→ xformers (단, QwenImageLayered는 dual-stream으로 호환 안될 수 있음)")
    elif mem_efficient_available:
        print("→ Memory Efficient SDPA")
    elif flash_available:
        print("→ Flash SDPA")
    else:
        print("→ Math SDPA (fallback, 느림)")

    print("=" * 60)

    return {
        "bnb": bnb_available,
        "xformers": xformers_available,
        "sageattention": sageattention_available,
        "flash_sdp": flash_available,
        "mem_efficient_sdp": mem_efficient_available,
    }

# 시작 시 VRAM/RAM 정리
def clear_memory():
    """VRAM과 RAM 정리"""
    gc.collect()
    if torch.cuda.is_available():
        torch.cuda.empty_cache()
        torch.cuda.reset_peak_memory_stats()
        torch.cuda.synchronize()

# 환경 정보 출력
ENV_INFO = print_environment_info()
print()
print("메모리 정리 중...")
clear_memory()
print("메모리 정리 완료")

from diffusers import QwenImageLayeredPipeline
from PIL import Image
from pptx import Presentation
import os
import gradio as gr
import uuid
import numpy as np
import random
import tempfile
import zipfile
from psd_tools import PSDImage
from huggingface_hub import snapshot_download
from pathlib import Path

MAX_SEED = np.iinfo(np.int32).max

# 모델 설정
MODEL_ID = "Qwen/Qwen-Image-Layered"
LOCAL_MODEL_DIR = Path(__file__).parent.parent / "models" / "Qwen-Image-Layered"

# RTX 4090 최적화 설정
def setup_cuda_optimizations(enable_flash_attention=False):
    """RTX 4090 CUDA 최적화

    Args:
        enable_flash_attention: Flash Attention 활성화 여부 (Windows에서 불안정할 수 있음)
    """
    # TF32 활성화 (Ada/Ampere GPU에서 성능 향상)
    torch.backends.cuda.matmul.allow_tf32 = True
    torch.backends.cudnn.allow_tf32 = True
    print("  - TF32 활성화")

    # cuDNN 벤치마크 모드 (입력 크기가 일정할 때 최적의 알고리즘 선택)
    torch.backends.cudnn.benchmark = True
    print("  - cuDNN 벤치마크 모드 활성화")

    # Flash Attention / SDPA 설정
    if enable_flash_attention:
        try:
            torch.backends.cuda.enable_flash_sdp(True)
            torch.backends.cuda.enable_mem_efficient_sdp(True)
            print("  - Flash Attention / Memory Efficient SDPA 활성화")
        except Exception as e:
            print(f"  - Flash Attention 설정 실패 (무시됨): {e}")
    else:
        # 명시적으로 비활성화하여 기본 SDPA 사용
        try:
            torch.backends.cuda.enable_flash_sdp(False)
            torch.backends.cuda.enable_mem_efficient_sdp(True)  # Memory efficient는 안정적
            print("  - Flash Attention 비활성화, Memory Efficient SDPA만 사용")
        except Exception as e:
            print(f"  - SDPA 설정 스킵: {e}")

    print("RTX 4090 CUDA 최적화 적용 완료")

def get_model_path():
    """로컬에 모델이 있으면 로컬 경로 반환, 없으면 다운로드 후 경로 반환"""
    if LOCAL_MODEL_DIR.exists() and any(LOCAL_MODEL_DIR.iterdir()):
        print(f"로컬 모델 사용: {LOCAL_MODEL_DIR}")
        return str(LOCAL_MODEL_DIR)

    print(f"모델 다운로드 중: {MODEL_ID} -> {LOCAL_MODEL_DIR}")
    LOCAL_MODEL_DIR.mkdir(parents=True, exist_ok=True)
    snapshot_download(
        repo_id=MODEL_ID,
        local_dir=str(LOCAL_MODEL_DIR),
    )
    print(f"다운로드 완료: {LOCAL_MODEL_DIR}")
    return str(LOCAL_MODEL_DIR)

def load_pipeline():
    """Pipeline 로드 및 RTX 4090 최적화"""
    print("=" * 50)
    print("Pipeline 로딩 시작...")
    print("=" * 50)

    # CUDA 최적화 적용 (Flash Attention은 기본 비활성화)
    setup_cuda_optimizations(enable_flash_attention=False)

    model_path = get_model_path()
    print(f"모델 경로: {model_path}")

    # BitsAndBytes 양자화 설정 (QUANTIZATION_MODE 플래그에 따라)
    quantization_config = None
    if QUANTIZATION_MODE:
        try:
            from diffusers.quantizers import PipelineQuantizationConfig
            if QUANTIZATION_MODE == "4bit":
                # 4bit NF4 양자화 (메모리 75% 감소, 품질 저하 가능)
                quantization_config = PipelineQuantizationConfig(
                    quant_backend="bitsandbytes_4bit",
                    quant_kwargs={
                        "bnb_4bit_compute_dtype": torch.bfloat16,
                        "bnb_4bit_quant_type": "nf4",
                    },
                    components_to_quantize=["transformer", "text_encoder"],
                )
                print("  - BitsAndBytes 4bit NF4 양자화 (transformer + text_encoder)")
            elif QUANTIZATION_MODE == "4bit_quality":
                # 4bit NF4 양자화 - transformer만 (text_encoder는 원본 유지)
                quantization_config = PipelineQuantizationConfig(
                    quant_backend="bitsandbytes_4bit",
                    quant_kwargs={
                        "bnb_4bit_compute_dtype": torch.bfloat16,
                        "bnb_4bit_quant_type": "nf4",
                    },
                    components_to_quantize=["transformer"],  # text_encoder 제외
                )
                print("  - BitsAndBytes 4bit NF4 양자화 (transformer만, text_encoder 원본 유지)")
            elif QUANTIZATION_MODE == "8bit":
                # 8bit 양자화 (메모리 50% 감소, 품질 유지)
                quantization_config = PipelineQuantizationConfig(
                    quant_backend="bitsandbytes_8bit",
                    quant_kwargs={
                        "llm_int8_threshold": 6.0,  # 기본값
                    },
                    components_to_quantize=["transformer", "text_encoder"],
                )
                print("  - BitsAndBytes 8bit 양자화 (transformer + text_encoder)")
        except ImportError:
            print("  - PipelineQuantizationConfig 미지원 - 양자화 없이 로드")
        except Exception as e:
            print(f"  - 양자화 설정 실패: {e}")
    else:
        print("  - 양자화 비활성화 (bfloat16 사용)")

    print("모델 로딩 중...")
    if quantization_config:
        pipe = QwenImageLayeredPipeline.from_pretrained(
            model_path,
            quantization_config=quantization_config,
            torch_dtype=torch.bfloat16,
        )
    else:
        pipe = QwenImageLayeredPipeline.from_pretrained(
            model_path,
            torch_dtype=torch.bfloat16,
        )
    print("모델 로딩 완료")

    # GPU 전송 또는 CPU Offload 설정
    if CPU_OFFLOAD_MODE == "off":
        print("GPU로 전송 중...")
        pipe = pipe.to("cuda")
        print("GPU 전송 완료")
    elif CPU_OFFLOAD_MODE == "model":
        print("Model CPU Offload 활성화 중...")
        pipe.enable_model_cpu_offload()
        print("Model CPU Offload 활성화 완료 (VRAM 절약, 추론 시 GPU 사용)")
    elif CPU_OFFLOAD_MODE == "sequential":
        print("Sequential CPU Offload 활성화 중...")
        pipe.enable_sequential_cpu_offload()
        print("Sequential CPU Offload 활성화 완료 (최소 VRAM 사용)")
    else:
        print(f"알 수 없는 CPU_OFFLOAD_MODE: {CPU_OFFLOAD_MODE}, GPU로 전송...")
        pipe = pipe.to("cuda")

    pipe.set_progress_bar_config(disable=None)

    # VAE를 sliced/tiled 모드로 설정 (메모리 효율)
    if hasattr(pipe, 'vae'):
        pipe.vae.enable_slicing()
        pipe.vae.enable_tiling()
        print("  - VAE slicing/tiling 활성화")

    # FirstBlockCache 적용 (TeaCache와 유사 - 추론 속도 20-40% 향상)
    if ENABLE_FIRST_BLOCK_CACHE:
        try:
            from diffusers.hooks import apply_first_block_cache, FirstBlockCacheConfig
            apply_first_block_cache(pipe.transformer, FirstBlockCacheConfig(threshold=FIRST_BLOCK_CACHE_THRESHOLD))
            print(f"  - FirstBlockCache 활성화 (threshold={FIRST_BLOCK_CACHE_THRESHOLD})")
        except Exception as e:
            print(f"  - FirstBlockCache 적용 실패: {e}")
    else:
        print("  - FirstBlockCache 비활성화")

    # 이 모델은 커스텀 QwenDoubleStreamAttnProcessor2_0을 사용
    # xformers는 호환되지 않음 (dual-stream output 미지원)
    # Native PyTorch SDPA 사용 (PyTorch 2.6.0에서 최적화됨)
    print("  - Native SDPA 사용 (QwenDoubleStreamAttnProcessor2_0)")

    # torch.compile 비활성화 (안정성 우선)
    ENABLE_TORCH_COMPILE = False
    if ENABLE_TORCH_COMPILE:
        try:
            pipe.transformer = torch.compile(pipe.transformer, mode="reduce-overhead")
            print("  - torch.compile 적용됨 (mode=reduce-overhead)")
        except Exception as e:
            print(f"  - torch.compile 스킵: {e}")

    print("=" * 50)
    print("Pipeline 준비 완료!")
    print("=" * 50)
    return pipe

pipeline = load_pipeline()

def imagelist_to_pptx(img_files):
    with Image.open(img_files[0]) as img:
        img_width_px, img_height_px = img.size

    def px_to_emu(px, dpi=96):
        inch = px / dpi
        emu = inch * 914400
        return int(emu)

    prs = Presentation()
    prs.slide_width = px_to_emu(img_width_px)
    prs.slide_height = px_to_emu(img_height_px)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = top = 0
    for img_path in img_files:
        slide.shapes.add_picture(img_path, left, top, width=px_to_emu(img_width_px), height=px_to_emu(img_height_px))

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return tmp.name

def imagelist_to_psd(img_files):
    layers = []
    for path in img_files:
        layers.append(Image.open(path).convert('RGBA'))

    width, height = layers[0].size
    psd = PSDImage.new(mode='RGBA', size=(width, height))

    for i, img in enumerate(layers):
        name = f"Layer {i + 1}"
        layer = psd.create_pixel_layer(image=img, name=name)
        psd.append(layer)
    
    with tempfile.NamedTemporaryFile(suffix=".psd", delete=False) as tmp:
        psd.save(tmp.name)
        return tmp.name

def export_gallery(images):
    # images: list of image file paths
    images = [e[0] for e in images]
    pptx_path = imagelist_to_pptx(images)
    return pptx_path

def infer(input_image,
          seed=777,
          randomize_seed=False,
          prompt=None,
          neg_prompt=" ",
          true_guidance_scale=4.0,
          num_inference_steps=50,
          layer=4,
          cfg_norm=True,
          use_en_prompt=True):
    
    if randomize_seed:
        seed = random.randint(0, MAX_SEED)
        
    if isinstance(input_image, list):
        input_image = input_image[0]
        
    if isinstance(input_image, str):
        pil_image = Image.open(input_image).convert("RGB").convert("RGBA")
    elif isinstance(input_image, Image.Image):
        pil_image = input_image.convert("RGB").convert("RGBA")
    elif isinstance(input_image, np.ndarray):
        pil_image = Image.fromarray(input_image).convert("RGB").convert("RGBA")
    else:
        raise ValueError("Unsupported input_image type: %s" % type(input_image))
    
    # CPU Offload 모드에서는 generator 장치를 CPU로 설정
    if CPU_OFFLOAD_MODE in ("model", "sequential"):
        generator_device = 'cpu'
    else:
        generator_device = 'cuda'

    inputs = {
        "image": pil_image,
        "generator": torch.Generator(device=generator_device).manual_seed(seed),
        "true_cfg_scale": true_guidance_scale,
        "prompt": prompt,
        "negative_prompt": neg_prompt,
        "num_inference_steps": num_inference_steps,
        "num_images_per_prompt": 1,
        "layers": layer,
        "resolution": 640,      # Using different bucket (640, 1024) to determine the resolution. For this version, 640 is recommended
        "cfg_normalize": cfg_norm,  # Whether enable cfg normalization.
        "use_en_prompt": use_en_prompt,
    }
    print(inputs)
    with torch.inference_mode(), torch.amp.autocast('cuda', dtype=torch.bfloat16):
        output = pipeline(**inputs)
        output_images = output.images[0]

    # VRAM/RAM 정리
    clear_memory()
    
    output = []
    temp_files = []
    for i, image in enumerate(output_images):
        output.append(image)
        # Save to temp file for export
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        image.save(tmp.name)
        temp_files.append(tmp.name)
    
    # Generate PPTX
    pptx_path = imagelist_to_pptx(temp_files)
    
    # Generate ZIP
    with tempfile.NamedTemporaryFile(suffix=".zip", delete=False) as tmp:
        with zipfile.ZipFile(tmp.name, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for i, img_path in enumerate(temp_files):
                zipf.write(img_path, f"layer_{i+1}.png")
        zip_path = tmp.name
    
    # Generate PSD
    psd_path = imagelist_to_psd(temp_files)
    return output, pptx_path, zip_path, psd_path

examples = [
            "assets/test_images/1.png",
            "assets/test_images/2.png",
            "assets/test_images/3.png",
            "assets/test_images/4.png",
            "assets/test_images/5.png",
            "assets/test_images/6.png",
            "assets/test_images/7.png",
            "assets/test_images/8.png",
            "assets/test_images/9.png",
            "assets/test_images/10.png",
            "assets/test_images/11.png",
            "assets/test_images/12.png",
            "assets/test_images/13.png",
            ]


with gr.Blocks() as demo:
    with gr.Column(elem_id="col-container"):
        gr.HTML('''<p align="center"><img src="https://qianwen-res.oss-cn-beijing.aliyuncs.com/Qwen-Image/layered/qwen-image-layered-logo.png" width="400"/><p>''')
        gr.Markdown("""
                    The text prompt is intended to describe the overall content of the input image—including elements that may be partially occluded (e.g., you may specify the text hidden behind a foreground object). It is not designed to control the semantic content of individual layers explicitly.
                    """)
        with gr.Row():
            with gr.Column(scale=1):
                input_image = gr.Image(label="Input Image", image_mode="RGBA")
                
                
                with gr.Accordion("Advanced Settings", open=False):
                    prompt = gr.Textbox(
                        label="Prompt (Optional)",
                        placeholder="Please enter the prompt to describe the image. It is not designed to control the semantic content of individual layers explicitly. (Optional)",
                        value="",
                        lines=3,
                    )
                    neg_prompt = gr.Textbox(
                        label="Negative Prompt (Optional)",
                        placeholder="Please enter the negative prompt",
                        value=" ",
                        lines=3,
                    )
                    
                    seed = gr.Slider(
                        label="Seed",
                        minimum=0,
                        maximum=MAX_SEED,
                        step=1,
                        value=0,
                    )
                    randomize_seed = gr.Checkbox(label="Randomize seed", value=True)
                    
                    true_guidance_scale = gr.Slider(
                        label="True guidance scale",
                        minimum=1.0,
                        maximum=10.0,
                        step=0.1,
                        value=4.0
                    )

                    num_inference_steps = gr.Slider(
                        label="Number of inference steps",
                        minimum=1,
                        maximum=50,
                        step=1,
                        value=20,  # 20 steps로 감소 (4bit + FirstBlockCache와 함께 사용)
                    )

                    layer = gr.Slider(
                        label="Layers",
                        minimum=2,
                        maximum=10,
                        step=1,
                        value=4,
                    )

                    cfg_norm = gr.Checkbox(label="Whether enable CFG normalization", value=True)
                    use_en_prompt = gr.Checkbox(label="Automatic caption language if no prompt provided, True for EN, False for ZH", value=True)
                
                run_button = gr.Button("Decompose!", variant="primary")

            with gr.Column(scale=2):
                gallery = gr.Gallery(label="Layers", columns=4, rows=1, format="png")
                with gr.Row():
                    export_file = gr.File(label="Download PPTX")
                    export_zip_file = gr.File(label="Download ZIP")
                    export_psd_file = gr.File(label="Download PSD")

    gr.Examples(examples=examples,
                inputs=[input_image], 
                outputs=[gallery, export_file, export_zip_file, export_psd_file],
                fn=infer, 
                examples_per_page=14,
                cache_examples=False,
                run_on_click=True
    )

    run_button.click(
        fn=infer,
        inputs=[
            input_image,
            seed,
            randomize_seed,
            prompt,
            neg_prompt,
            true_guidance_scale,
            num_inference_steps,
            layer,
            cfg_norm,
            use_en_prompt,
        ], 
        outputs=[gallery, export_file, export_zip_file, export_psd_file],
    )

demo.launch(
    server_name="127.0.0.1",
    server_port=7869,
)
